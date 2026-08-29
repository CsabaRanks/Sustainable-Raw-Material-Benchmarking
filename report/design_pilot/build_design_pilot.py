# Design pilot: cover + main pages 1-3.
# Architecture per the approved report storyline (frozen).
import os
from pptx import Presentation
from pptx.util import Inches as In, Pt, Emu
from pptx.dml.color import RGBColor as C
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "Sustainable_Raw_Material_Benchmarking_Design_Pilot.pptx")

# ── palette ───────────────────────────────────────────────────────────────
INK      = C(0x14, 0x18, 0x1B)
INK2     = C(0x3D, 0x47, 0x4D)
INK3     = C(0x6E, 0x7B, 0x82)
RULE     = C(0xD4, 0xDA, 0xDE)
RULE_LT  = C(0xE8, 0xEC, 0xEF)
PAPER    = C(0xFF, 0xFF, 0xFF)
OFFWHITE = C(0xF7, 0xF9, 0xFA)
BLUE     = C(0x2E, 0x5A, 0x87)
BLUE_LT  = C(0xBF, 0xD2, 0xE4)
BLUE_XLT = C(0xE4, 0xED, 0xF5)
GREEN    = C(0x3E, 0x8F, 0x63)
GREEN_LT = C(0xD9, 0xEA, 0xE0)
RED      = C(0xB3, 0x37, 0x2C)
RED_LT   = C(0xF4, 0xDF, 0xDC)
AMBER    = C(0xD9, 0x9A, 0x2B)
AMBER_LT = C(0xFA, 0xEC, 0xD4)

F_L, F_SL, F_R, F_SB = "Segoe UI Light", "Segoe UI Semilight", "Segoe UI", "Segoe UI Semibold"

# ── grid ──────────────────────────────────────────────────────────────────
W, H = 13.333, 7.5
ML, MR = 0.78, 0.78
CW = W - ML - MR                      # 11.773
Y_EYE, Y_TITLE, Y_LEAD, Y_RULE, Y_EX, Y_FOOT = 0.46, 0.72, 1.80, 2.42, 2.66, 6.84
EX_H = 6.60 - Y_EX

prs = Presentation()
prs.slide_width, prs.slide_height = In(W), In(H)
BLANK = prs.slide_layouts[6]


def txt(sl, x, y, w, h, runs, size=10, font=F_R, color=INK, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, spacing=1.0, space_after=0, caps=False, wrap=True):
    """runs: str, or list of (text, {overrides}) or list of paragraphs [[...],[...]]"""
    tb = sl.shapes.add_textbox(In(x), In(y), In(w), In(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    paras = runs if isinstance(runs, list) and runs and isinstance(runs[0], list) else [runs]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        p.space_after = Pt(space_after)
        items = para if isinstance(para, list) else [(para, {})]
        for t, ov in items:
            r = p.add_run()
            r.text = t.upper() if caps else t
            f = r.font
            f.name = ov.get("font", font)
            f.size = Pt(ov.get("size", size))
            f.color.rgb = ov.get("color", color)
            f.bold = ov.get("bold", False)
            f.italic = ov.get("italic", False)
    return tb


def rect(sl, x, y, w, h, fill=None, line=None, lw=0.75, shape=MSO_SHAPE.RECTANGLE):
    s = sl.shapes.add_shape(shape, In(x), In(y), In(w), In(h))
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(lw)
    if s.has_text_frame:
        s.text_frame.word_wrap = True
        s.text_frame.margin_left = s.text_frame.margin_right = 0
        s.text_frame.margin_top = s.text_frame.margin_bottom = 0
    return s


def hline(sl, x, y, w, color=RULE, lw=0.75):
    c = sl.shapes.add_connector(1, In(x), In(y), In(x + w), In(y))
    c.line.color.rgb = color; c.line.width = Pt(lw)
    return c


ACTS = ["I  THE MANAGEMENT PROBLEM", "II  WHAT BEST PRACTICE LOOKS LIKE",
        "III  FROM PRINCIPLE TO ANALYTICS", "IV  BEYOND CARBON", "V  MANAGEMENT IMPLICATIONS"]


def page_frame(sl, act_idx, eyebrow, title_lines, lead, source, pageno):
    txt(sl, ML, Y_EYE, 7.4, 0.22, eyebrow, size=8.5, font=F_SB, color=INK3, caps=True)
    # act progress: five segments, current filled
    sx = W - MR - (5 * 0.30 + 4 * 0.07)
    for i in range(5):
        rect(sl, sx + i * 0.37, Y_EYE + 0.045, 0.30, 0.075,
             fill=BLUE if i == act_idx else RULE_LT)
    txt(sl, ML, Y_TITLE, CW - 0.4, 1.0, title_lines, size=22, font=F_SB,
        color=INK, spacing=1.06)
    if lead:
        txt(sl, ML, Y_LEAD, CW - 2.3, 0.5, lead, size=11, font=F_R, color=INK2, spacing=1.24)
    hline(sl, ML, Y_RULE, CW, RULE, 0.75)
    txt(sl, ML, Y_FOOT, CW - 1.0, 0.42, source, size=7.5, font=F_R, color=INK3, spacing=1.22)
    txt(sl, W - MR - 0.6, Y_FOOT, 0.6, 0.22, str(pageno), size=9, font=F_SB,
        color=INK3, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════ COVER
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=PAPER)
rect(s, 0, 0, 0.13, H, fill=INK)                      # spine bar

LX = 1.30
txt(s, LX, 1.02, 6.0, 0.22, "BENCHMARK REPORT   ·   2026", size=8.5, font=F_SB,
    color=INK3, caps=True)
hline(s, LX, 1.42, 1.5, INK, 1.4)

txt(s, LX, 1.86, 7.2, 2.0,
    [[("SUSTAINABLE", {})], [("RAW MATERIAL", {})], [("BENCHMARKING", {})]],
    size=40, font=F_L, color=INK, spacing=1.02)

txt(s, LX, 4.06, 6.6, 0.4, "Best Practices, Decision Framework & Analytics Prototype",
    size=15.5, font=F_SL, color=BLUE)

# synthetic-data disclosure — bordered, elegant
rect(s, LX, 4.86, 6.15, 0.86, fill=OFFWHITE)
rect(s, LX, 4.86, 0.035, 0.86, fill=BLUE)
txt(s, LX + 0.26, 5.04, 5.65, 0.6,
    [[("A decision method demonstrated on synthetic data.", {"font": F_SB, "color": INK})],
     [("No supplier, company or market data is used or implied.", {"color": INK2})]],
    size=9.5, spacing=1.32)

txt(s, LX, 6.28, 6.6, 0.25,
    [[("ELIGIBILITY", {"font": F_SB, "color": INK}), ("     ·     ", {"color": RULE}),
      ("COMPARABILITY", {"font": F_SB, "color": INK}), ("     ·     ", {"color": RULE}),
      ("TRADE-OFF", {"font": F_SB, "color": INK})]], size=9, caps=False)
hline(s, LX, 6.16, 6.15, RULE, 0.75)

# cover mark: 5 cases x 5 materials, tinted by the dataset's cost-delta pattern
sign = [
    [0, -3, -1, 2, 1],
    [0, -3, 3, -1, 2],
    [0, 1, -2, 3, -1],
    [0, -3, -1, -2, 1],
    [0, -1, 1, 0, 0],
]
TINT = {-3: C(0x7F, 0xB8, 0x99), -2: C(0xA8, 0xCE, 0xB8), -1: C(0xD2, 0xE7, 0xDB),
        0: C(0xEC, 0xEF, 0xF1), 1: C(0xF2, 0xD8, 0xD4), 2: C(0xE6, 0xB4, 0xAC),
        3: C(0xD4, 0x8B, 0x80)}
GX, GY, CS, GP = 9.45, 1.86, 0.44, 0.115
for r in range(5):
    for c_ in range(5):
        rect(s, GX + c_ * (CS + GP), GY + r * (CS + GP), CS, CS, fill=TINT[sign[r][c_]])
txt(s, GX, GY + 5 * (CS + GP) + 0.18, 3.0, 0.5,
    [[("25 material alternatives  ·  five benchmark cases", {"font": F_SB, "color": INK2})],
     [("green improves on cost, red deteriorates  ·  synthetic", {"color": INK3})]],
    size=8, spacing=1.34)

txt(s, LX, 6.72, 6.6, 0.22,
    "Sustainable Raw Material Benchmarking  ·  Design pilot  ·  v0.1  ·  Csaba Bakay",
    size=8, font=F_R, color=INK3)

# ══════════════════════════════════════════════════════════ PAGE 1
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=PAPER)
page_frame(
    s, 0, ACTS[0],
    [[("No raw material is universally sustainable — a defensible decision", {})],
     [("proves eligibility, then comparability, then shows the trade-off", {})]],
    "Sourcing decisions fail in three specific ways. The answer is not a better score; "
    "it is a sequence that makes each failure mode visible before the next one is judged.",
    "GHG Protocol Product Life Cycle Accounting and Reporting Standard (WRI/WBCSD, 2011), §1.5 and Appendix A  ·  "
    "European Commission JRC, Environmental and economic assessment of plastic waste recycling, JRC132067 (2023).  "
    "The three failure modes are a project synthesis of these sources.", 1)

PW, PG = 3.545, 0.569
for i, (num, name, q, fail, ctrl) in enumerate([
    ("01", "ELIGIBILITY", "Can this alternative be considered at all?",
     "An option is compared that cannot legally or technically be bought.",
     "Screen on true constraints first — and keep what was screened out visible."),
    ("02", "COMPARABILITY", "Were these numbers ever comparable?",
     "Two figures are compared that were never calculated on the same basis.",
     "Prove equivalent function, boundary, allocation, period and assurance."),
    ("03", "TRADE-OFF", "What are we giving up, and who owns that call?",
     "One dimension improves while another quietly moves the wrong way.",
     "Hold the dimensions on separate axes and name the decision owner."),
]):
    x = ML + i * (PW + PG)
    rect(s, x, Y_EX, PW, 3.16, fill=OFFWHITE)
    rect(s, x, Y_EX, PW, 0.035, fill=BLUE)
    txt(s, x + 0.30, Y_EX + 0.30, 0.9, 0.3, num, size=17, font=F_L, color=BLUE_LT)
    txt(s, x + 0.30, Y_EX + 0.72, PW - 0.6, 0.24, name, size=12, font=F_SB, color=INK)
    txt(s, x + 0.30, Y_EX + 1.02, PW - 0.6, 0.4, q, size=10, font=F_SL, color=INK2, spacing=1.2)
    hline(s, x + 0.30, Y_EX + 1.56, PW - 0.6, RULE, 0.75)
    rect(s, x + 0.30, Y_EX + 1.76, 0.075, 0.075, fill=RED)
    txt(s, x + 0.48, Y_EX + 1.70, PW - 0.78, 0.5,
        [[("Failure   ", {"font": F_SB, "color": RED, "size": 8}), (fail, {})]],
        size=9, color=INK2, spacing=1.24)
    rect(s, x + 0.30, Y_EX + 2.42, 0.075, 0.075, fill=GREEN)
    txt(s, x + 0.48, Y_EX + 2.36, PW - 0.78, 0.6,
        [[("Control   ", {"font": F_SB, "color": GREEN, "size": 8}), (ctrl, {})]],
        size=9, color=INK2, spacing=1.24)
    if i < 2:
        txt(s, x + PW + 0.10, Y_EX + 1.36, 0.4, 0.34, "\u203A", size=24, font=F_L,
            color=C(0xA3, 0xAF, 0xB6), align=PP_ALIGN.CENTER)

txt(s, ML, Y_EX + 3.40, CW, 0.24,
    [[("Each control is demonstrated later on the prototype: ", {"color": INK3}),
      ("eligibility on page 10", {"font": F_SB, "color": INK2}), ("  ·  ", {"color": RULE}),
      ("comparability on page 11", {"font": F_SB, "color": INK2}), ("  ·  ", {"color": RULE}),
      ("trade-off on pages 9 and 13", {"font": F_SB, "color": INK2})]], size=8.5)

# ══════════════════════════════════════════════════════════ PAGE 2
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=PAPER)
page_frame(
    s, 0, ACTS[0],
    [[("Carbon is the dimension we measure best — and the carbon", {})],
     [("standard itself says it is not the whole answer", {})]],
    "A product carbon footprint is decision-relevant and, in the chemical context, unusually "
    "well supported by supplier-exchange infrastructure. It is still one impact category.",
    "Quotation: GHG Protocol Product Life Cycle Accounting and Reporting Standard (WRI/WBCSD, 2011), §1.7, read in the original; emphasis added.  ·  "
    "Methods for other dimensions: ISO 14046:2014 (product water footprint); Commission Recommendation (EU) 2021/2279, Product Environmental "
    "Footprint, 16 impact categories.  ·  The data-availability constraint is a source-supported interpretation.", 2)

# pull quote
rect(s, ML, Y_EX, CW, 1.30, fill=OFFWHITE)
rect(s, ML, Y_EX, 0.035, 1.30, fill=BLUE)
txt(s, ML + 0.34, Y_EX + 0.22, CW - 0.8, 0.7,
    [[("“The limitation of a GHG-only inventory is that potential trade-offs or co-benefits "
       "between environmental impacts can be missed. Therefore, ", {}),
      ("the results of a GHG-only inventory should not be used to communicate the overall "
       "environmental performance of a product.", {"font": F_SB}),
      ("”", {})]],
    size=13, font=F_SL, color=INK, spacing=1.30)
txt(s, ML + 0.34, Y_EX + 1.00, CW - 0.8, 0.2,
    "GHG Protocol Product Standard, §1.7 \u2014 Limitations of product GHG inventories",
    size=8.5, font=F_SB, color=INK3)

# two columns
CY, CH = Y_EX + 1.62, 1.52
CWD = (CW - 0.55) / 2
txt(s, ML, CY, CWD, 0.22, "What the number covers", size=10.5, font=F_SB, color=INK)
hline(s, ML, CY + 0.30, CWD, INK, 1.0)
for j, (lbl, sub) in enumerate([("Climate change", "one impact category"),
                                ("Cradle-to-gate", "a partial life cycle"),
                                ("kg CO\u2082e per kg", "a declared unit, not a function")]):
    y = CY + 0.46 + j * 0.36
    rect(s, ML, y + 0.055, 0.075, 0.075, fill=BLUE)
    txt(s, ML + 0.20, y, CWD - 0.25, 0.3,
        [[(lbl, {"font": F_SB, "color": INK}), ("   " + sub, {"color": INK3, "size": 9})]], size=10)

X2 = ML + CWD + 0.55
txt(s, X2, CY, CWD, 0.22, "What it is silent on \u2014 named by the standard itself",
    size=10.5, font=F_SB, color=INK)
hline(s, X2, CY + 0.30, CWD, INK, 1.0)
for j, lbl in enumerate(["Ecosystem degradation", "Resource depletion",
                         "Ozone depletion", "Negative human health impacts"]):
    y = CY + 0.44 + j * 0.27
    rect(s, X2, y + 0.05, 0.075, 0.075, fill=RULE)
    txt(s, X2 + 0.20, y, CWD - 0.25, 0.24, lbl, size=10, color=INK2)

# methods strip
SY = CY + CH + 0.28
hline(s, ML, SY, CW, RULE, 0.75)
txt(s, ML, SY + 0.16, 3.5, 0.24, "And methods exist for the rest", size=10.5, font=F_SB, color=INK)
for j, (t, sub) in enumerate([("ISO 14046", "product water footprint"),
                              ("EU PEF", "16 impact categories")]):
    x = ML + 3.60 + j * 2.35
    rect(s, x, SY + 0.13, 2.15, 0.34, fill=BLUE_XLT)
    txt(s, x + 0.14, SY + 0.185, 1.9, 0.26,
        [[(t, {"font": F_SB, "color": BLUE}), ("  " + sub, {"color": INK2, "size": 8.5})]], size=9.5)
txt(s, ML + 8.45, SY + 0.15, CW - 8.45, 0.4,
    "The binding constraint is comparable supplier data, not the absence of a method.",
    size=9, font=F_SL, color=INK2, spacing=1.2)

# ══════════════════════════════════════════════════════════ PAGE 3
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=PAPER)
page_frame(
    s, 0, ACTS[0],
    [[("Which dimensions matter depends on the material, its feedstock", {})],
     [("and its origin — not on a universal checklist", {})]],
    "A universal ESG data request produces effort, not decisions. The constant is the screening "
    "question; the answer changes with what is bought, what it is made from and where it comes from.",
    "Classification and drivers: author synthesis, built on ISO 14046:2014, Commission Recommendation (EU) 2021/2279 "
    "(PEFCRs fix relevant impact categories per product group), Regulation (EU) 2024/1252 and Directive (EU) 2026/470.  ·  "
    "The classification is a project synthesis on sourced dimension definitions.  ·  Supplier-level social performance data is class E \u2014 currently insufficiently measurable.", 3)

# screening question band
rect(s, ML, Y_EX, CW, 0.40, fill=INK)
txt(s, ML + 0.28, Y_EX + 0.085, CW - 0.6, 0.26,
    [[("The screening question   ", {"font": F_SB, "color": C(0x8F, 0xA6, 0xB6), "size": 9}),
      ("Could this dimension materially change the decision for this material, "
       "from this feedstock, from this origin?", {"color": PAPER})]], size=10.5)

# matrix
TY = Y_EX + 0.60
LBLW, CLSW, DRVW = 4.35, 0.90, 1.58
DRIVERS = ["MATERIAL", "FEEDSTOCK", "ORIGIN", "APPLICATION"]
txt(s, ML, TY, LBLW, 0.2, "Decision dimension", size=8, font=F_SB, color=INK3, caps=True)
txt(s, ML + LBLW, TY, CLSW, 0.2, "Class", size=8, font=F_SB, color=INK3, caps=True,
    align=PP_ALIGN.CENTER)
for j, d in enumerate(DRIVERS):
    txt(s, ML + LBLW + CLSW + j * DRVW, TY, DRVW, 0.2, d, size=8, font=F_SB,
        color=INK3, align=PP_ALIGN.CENTER)
hline(s, ML, TY + 0.26, CW, INK, 1.0)

CLS_STYLE = {"A": (BLUE, PAPER), "B": (BLUE_LT, INK), "C": (INK, PAPER),
             "D": (RULE_LT, INK3), "E": (OFFWHITE, INK3)}
ROWS = [
    ("Climate / GHG \u2014 product carbon footprint", "A", []),
    ("Technical performance and qualification",       "C", [0, 3]),
    ("Specific legal compliance", "C", [0, 2, 3]),
    ("Commercial performance \u2014 price, TCO",       "A", [0, 3]),
    ("Circularity / resource loops",                   "B", [0, 1, 3]),
    ("Pollution / toxicity",                           "B", [0, 3]),
    ("Water",                                          "B", [1, 2]),
    ("Resource use",                                   "B", [0, 1]),
    ("Supply-chain resilience",                        "B", [0, 2]),
    ("Land use / biodiversity",                        "D", [1, 2]),
]
RH = 0.27
for i, (name, cls, drv) in enumerate(ROWS):
    y = TY + 0.34 + i * RH
    if i % 2 == 1:
        rect(s, ML, y - 0.045, CW, RH, fill=OFFWHITE)
    txt(s, ML + 0.02, y + 0.025, LBLW - 0.1, 0.24, name, size=9.5, color=INK)
    fc, tc = CLS_STYLE[cls]
    rect(s, ML + LBLW + CLSW / 2 - 0.115, y + 0.015, 0.23, 0.20, fill=fc)
    txt(s, ML + LBLW + CLSW / 2 - 0.115, y + 0.038, 0.23, 0.18, cls, size=8,
        font=F_SB, color=tc, align=PP_ALIGN.CENTER)
    for j in range(4):
        cx = ML + LBLW + CLSW + j * DRVW + DRVW / 2
        if j in drv:
            rect(s, cx - 0.055, y + 0.055, 0.11, 0.11, fill=BLUE, shape=MSO_SHAPE.OVAL)
        else:
            rect(s, cx - 0.04, y + 0.07, 0.08, 0.08, fill=RULE_LT, shape=MSO_SHAPE.OVAL)
hline(s, ML, TY + 0.34 + len(ROWS) * RH + 0.02, CW, RULE, 0.75)

# legend
LY = TY + 0.34 + len(ROWS) * RH + 0.16
for dx, cls, lbl in [(0.00, "A", "core"), (1.32, "B", "context-dependent"),
                     (3.32, "C", "gate / compliance"),
                     (5.32, "D", "big picture, not this prototype")]:
    fc, tc = CLS_STYLE[cls]
    rect(s, ML + dx, LY + 0.02, 0.19, 0.165, fill=fc)
    txt(s, ML + dx + 0.235, LY, 2.6, 0.2,
        [[(cls, {"font": F_SB, "color": INK}), ("  " + lbl, {"color": INK3})]], size=8.5)
rect(s, ML + 8.57, LY + 0.055, 0.11, 0.11, fill=BLUE, shape=MSO_SHAPE.OVAL)
txt(s, ML + 8.80, LY, 2.60, 0.2, "driver that can change materiality", size=8.5, color=INK3)

# ══════════════════════════════════════════════════════════ PAGE 4
# Mode: conceptual architecture. Four decision roles as four behaviours.
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=PAPER)
page_frame(
    s, 1, ACTS[1],
    [[("Most things called “gates” are not gates — and mixing the four", {})],
     [("decision roles is what breaks a benchmark", {})]],
    "Four dimensions can look alike on a data sheet and behave completely differently in a decision. "
    "A gate is not a preference, and evidence is not performance.",
    "Decision-role convention: project synthesis  ·  C6 Ashby, constraints vs. objectives  ·  "
    "A2 ISO 14044 and P3 PACT v3 for evidence metadata  ·  V2 REACH Annex XIV / XVII instruments  ·  V1 Directive (EU) 2026/470, applying from "
    "26 July 2029 to a narrow set of very large companies.  The four-role convention is a PROJECT SYNTHESIS; the legal statements are fact.", 4)

LW = 6.85                                    # left column: the four bands
RX, RW = ML + 7.15, 4.62                     # right column: two insets

# discreet but explicit synthesis label
rect(s, ML, Y_EX, LW, 0.24, fill=None, line=RULE, lw=0.75)
txt(s, ML + 0.14, Y_EX + 0.052, LW - 0.3, 0.2,
    [[("PROJECT SYNTHESIS", {"font": F_SB, "color": INK}),
      ("   a practical decision convention, not a validated external framework", {"color": INK3})]],
    size=8)

BY, BH, BG = Y_EX + 0.32, 0.78, 0.10
for i, (role, q, behav, key) in enumerate([
    ("GATE", "Can this alternative be considered at all?",
     "Pass / fail.\nNever traded off.", INK),
    ("PERFORMANCE", "How does an eligible alternative perform?",
     "Compared, plotted,\ntraded off.", BLUE),
    ("EVIDENCE", "How much confidence does that comparison deserve?",
     "Reported beside the number,\nnever blended into it.", AMBER),
    ("CONTEXT", "Under what material, application and supply conditions?",
     "Filters and explains.\nNot optimised.", RULE),
]):
    y = BY + i * (BH + BG)
    rect(s, ML, y, LW, BH, fill=OFFWHITE)
    rect(s, ML, y, 0.055, BH, fill=key)
    txt(s, ML + 0.26, y + 0.11, 4.0, 0.24, role, size=12, font=F_SB, color=INK)
    txt(s, ML + 0.26, y + 0.41, 4.05, 0.3, q, size=9.5, color=INK2)
    txt(s, ML + 4.52, y + 0.20, 2.15, 0.44,
        [[(behav.split("\n")[0], {})], [(behav.split("\n")[1], {})]],
        size=9, color=INK3, spacing=1.22)

txt(s, ML, BY + 4 * (BH + BG) + 0.04, LW, 0.22,
    [[("A gate without a named threshold owner is an opinion.", {"font": F_SB, "color": INK}),
      ("   Every gate names its threshold owner, decision owner and time horizon.", {"color": INK3})]],
    size=8.5)

# inset 1 — the role is set by the rule
txt(s, RX, Y_EX, RW, 0.2, "The role is set by the decision rule, not the topic name",
    size=8.5, font=F_SB, color=INK3, caps=True)
hline(s, RX, Y_EX + 0.26, RW, INK, 1.0)
for j, (topic, rule) in enumerate([
    ("Supply risk", "gate where a minimum volume or dual-source policy is contractual; "
                    "performance where lower concentration is preferred; context elsewhere"),
    ("Regulatory criteria", "gate only where the specific use is unlawful, unauthorised or "
                            "outside written company policy — disclosure duties are not gates"),
    ("Social criteria", "gate for prohibited conduct or an explicit minimum requirement; "
                        "general due diligence is a risk-based process, not a certificate"),
]):
    y = Y_EX + 0.36 + j * 0.50
    txt(s, RX, y, RW, 0.44,
        [[(topic + "  —  ", {"font": F_SB, "color": INK}), (rule, {"color": INK2})]],
        size=8.5, spacing=1.24)

# inset 2 — what is, and is not, a gate
IY = Y_EX + 1.94
txt(s, RX, IY, RW, 0.2, "What is — and is not — a gate", size=8.5, font=F_SB,
    color=INK3, caps=True)
hline(s, RX, IY + 0.26, RW, INK, 1.0)
txt(s, RX, IY + 0.36, RW, 0.2, "Not a gate", size=9, font=F_SB, color=RED)
for j, (t, sub) in enumerate([
    ("REACH Candidate-List presence", "a disclosure duty"),
    ("General due diligence status", "a risk-based process"),
]):
    y = IY + 0.60 + j * 0.24
    rect(s, RX + 0.02, y + 0.055, 0.075, 0.075, fill=RED)
    txt(s, RX + 0.20, y, RW - 0.25, 0.22,
        [[(t, {"color": INK2}), ("   " + sub, {"color": INK3})]], size=8.5)
txt(s, RX, IY + 1.14, RW, 0.2, "A gate", size=9, font=F_SB, color=INK)
for j, t in enumerate([
    "Annex XIV after sunset, without authorisation",
    "Annex XVII restriction covering the specific use",
    "An explicit company policy threshold",
]):
    y = IY + 1.38 + j * 0.24
    rect(s, RX + 0.02, y + 0.055, 0.075, 0.075, fill=INK)
    txt(s, RX + 0.20, y, RW - 0.25, 0.22, t, size=8.5, color=INK2)

# ══════════════════════════════════════════════════════════ PAGE 5
# Mode: best-practice sequence. The report's best-practice anchor.
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=PAPER)
page_frame(
    s, 1, ACTS[1] + "   ·   BEST-PRACTICE ANCHOR",
    [[("Screen, compare, qualify, decide, record — the practices that", {})],
     [("survive the evidence, and how strongly each is supported", {})]],
    "Thirteen practices came out of the research. Four are strongly evidenced. Showing which is which "
    "is itself a practice — it stops a plausible recommendation being read as an established one.",
    "BP-1 to BP-13: author synthesis, with the evidence recorded for each practice  ·  P1 GHG Protocol Product Standard §A.1  ·  P2 JRC132067  ·  "
    "P3 PACT v3  ·  A1 ISO 20400  ·  A2 ISO 14044.  The marker describes how well a practice is supported — never organisational maturity. "
    "No scoring, no ranking, no maturity model.", 5)

STEPS = [
    ("SCREEN", "eligibility", [
        ("BP-1", "Screen on hard constraints first", 1),
        ("BP-9", "Gate on specific legal thresholds, not general risk", 1),
        ("BP-8", "Treat supply risk as a dimension, not a footnote", 1)]),
    ("COMPARE", "comparability", [
        ("BP-13", "Define equivalent function before any per-kg benchmark", 2),
        ("BP-5", "Satisfy the six comparability conditions first", 2),
        ("BP-2", "Require the primary-data share, not just a figure", 1),
        ("BP-12", "Declare a comparability verdict, not a colour", 1)]),
    ("QUALIFY", "evidence", [
        ("BP-3", "Carry evidence quality beside the number, never inside it", 2),
        ("BP-4", "Never substitute zero for missing data", 2)]),
    ("DECIDE", "trade-off", [
        ("BP-6", "Decide on total cost, not on price per kg", 1),
        ("BP-10", "Price in regulatory cost transmission early", 1),
        ("BP-7", "Validate circularity against life-cycle impact", 1)]),
    ("RECORD", "owned", [
        ("BP-11", "One decision record, not two", 0)]),
]
SW, SG = (CW - 4 * 0.22) / 5, 0.22
for i, (name, spine, items) in enumerate(STEPS):
    x = ML + i * (SW + SG)
    rect(s, x, Y_EX, SW, 0.44, fill=BLUE)
    txt(s, x + 0.16, Y_EX + 0.055, SW - 0.3, 0.22, name, size=11, font=F_SB, color=PAPER)
    txt(s, x + 0.16, Y_EX + 0.26, SW - 0.3, 0.18, spine, size=8, color=BLUE_LT)
    if i < 4:
        txt(s, x + SW + 0.015, Y_EX + 0.10, 0.19, 0.26, "›", size=15, font=F_L,
            color=BLUE, align=PP_ALIGN.CENTER)
    for j, (bp, text, strength) in enumerate(items):
        y = Y_EX + 0.56 + j * 0.58
        if strength == 2:
            rect(s, x, y - 0.04, SW, 0.50, fill=BLUE_XLT)
        mx, my = x + 0.09, y + 0.045
        if strength == 2:
            rect(s, mx, my, 0.10, 0.10, fill=INK)
        elif strength == 1:
            rect(s, mx, my, 0.10, 0.10, fill=INK3)
        else:
            rect(s, mx, my, 0.10, 0.10, fill=PAPER, line=INK3, lw=0.75)
        txt(s, x + 0.26, y, SW - 0.36, 0.18, bp, size=7.5, font=F_SB,
            color=INK if strength == 2 else INK3)
        txt(s, x + 0.09, y + 0.20, SW - 0.18, 0.34, text, size=9,
            font=F_SB if strength == 2 else F_R,
            color=INK if strength == 2 else INK2, spacing=1.18)

txt(s, ML + 4 * (SW + SG) + 0.09, Y_EX + 1.16, SW - 0.18, 0.5,
    "Shared by procurement and engineering, with the gate reasons visible.",
    size=8.5, color=INK3, spacing=1.20)

LGY = Y_EX + 3.02
for dx, lbl, kind in [(0.00, "strongly evidenced — primary standard or EU authority", 2),
                      (3.60, "supported — cited, but narrower", 1),
                      (6.10, "project recommendation — plausible, not externally evidenced", 0)]:
    if kind == 2:
        rect(s, ML + dx, LGY + 0.035, 0.10, 0.10, fill=INK)
    elif kind == 1:
        rect(s, ML + dx, LGY + 0.035, 0.10, 0.10, fill=INK3)
    else:
        rect(s, ML + dx, LGY + 0.035, 0.10, 0.10, fill=PAPER, line=INK3, lw=0.75)
    txt(s, ML + dx + 0.20, LGY, 3.5, 0.2, lbl, size=8.5, color=INK3)

BY2 = Y_EX + 3.30
rect(s, ML, BY2, CW, 0.56, fill=INK)
txt(s, ML + 0.28, BY2 + 0.085, CW - 0.6, 0.42,
    [[("What a procurement organisation does differently", {"font": F_SB, "color": PAPER})],
     [("Name the rulebook in the RFQ  ·  report the primary-data share  ·  keep evidence beside the number  ·  "
       "never zero-fill a gap  ·  one decision record, not two", {"color": C(0xC6, 0xD4, 0xE0)})]],
    size=9.5, spacing=1.30)

# ══════════════════════════════════════════════════════════ PAGE 6
# Mode: comparability test.
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=PAPER)
page_frame(
    s, 1, ACTS[1],
    [[("Two supplier carbon footprints are not comparable until six", {})],
     [("conditions hold — and a PCF alone satisfies none of them", {})]],
    "The standard that defines the product carbon footprint also states that it does not, on its own, "
    "support a purchasing comparison. Comparability is engineered upstream in the specification.",
    "All six conditions verbatim from the GHG Protocol Product Life Cycle Accounting and Reporting Standard (WRI/WBCSD, 2011), Appendix A §A.1; "
    "Table A.2 states that business purchasing decisions are not supported by the standard unaided.  ·  D1 BASF disclosure of economic allocation "
    "where physical data is unavailable.  ·  B1 TfS, B3 Catena-X and P3 PACT v3 are examples of the additional programme specifications §A.1 "
    "requires; conformance to them was not verified here.  The three-state verdict is a source-supported interpretation.", 6)

CDW, CTR = 3.90, 3.573
for x, lbl in [(ML, "SUPPLIER A"), (ML + 7.87, "SUPPLIER B")]:
    rect(s, x, Y_EX, CDW, 0.72, fill=OFFWHITE)
    txt(s, x + 0.28, Y_EX + 0.14, CDW - 0.5, 0.22, lbl, size=10.5, font=F_SB, color=INK)
    txt(s, x + 0.28, Y_EX + 0.40, CDW - 0.5, 0.2,
        "Product carbon footprint, kg CO₂e per kg", size=9, color=INK3)
CX = ML + 4.10
# drawn, not typeset: glyph metrics cannot place the strike reliably
ECX = CX + CTR / 2
rect(s, ECX - 0.26, Y_EX + 0.255, 0.52, 0.055, fill=INK3)
rect(s, ECX - 0.26, Y_EX + 0.425, 0.52, 0.055, fill=INK3)
c = s.shapes.add_connector(1, In(ECX - 0.36), In(Y_EX + 0.3675),
                           In(ECX + 0.36), In(Y_EX + 0.3675))
c.line.color.rgb = RED; c.line.width = Pt(2.5)
txt(s, CX, Y_EX + 0.80, CTR, 0.22, "A PCF alone satisfies none of the six conditions.",
    size=9, font=F_SB, color=INK, align=PP_ALIGN.CENTER)

HY = Y_EX + 1.06
txt(s, ML, HY, 7.0, 0.2, "The six conditions for a business purchasing comparison",
    size=8.5, font=F_SB, color=INK3, caps=True)
txt(s, ML + 7.0, HY, CW - 7.0, 0.2,
    "GHG Protocol Product Standard, Appendix A §A.1 — verbatim",
    size=8.5, color=INK3, align=PP_ALIGN.RIGHT)
hline(s, ML, HY + 0.26, CW, INK, 1.0)

CONDS = [
    ("1", "The unit of analysis is identical", None),
    ("2", "System boundaries and the temporal boundary are equivalent", None),
    ("3", "The same allocation methods are used for similar processes",
     "Broken between two honest suppliers wherever one allocates economically and the other physically."),
    ("4", "Data types, data quality and uncertainty are reported and assessed", None),
    ("5", "Temporal and geographical representativeness is assessed", None),
    ("6", "Third-party assurance has been performed",
     "The least met of the six in practice."),
]
COLW = (CW - 0.50) / 2
for k, (num, text, note) in enumerate(CONDS):
    col, row = k % 2, k // 2
    x = ML + col * (COLW + 0.50)
    y = HY + 0.38 + row * 0.54
    rect(s, x, y + 0.015, 0.22, 0.20, fill=BLUE)
    txt(s, x, y + 0.038, 0.22, 0.18, num, size=8, font=F_SB, color=PAPER, align=PP_ALIGN.CENTER)
    txt(s, x + 0.32, y, COLW - 0.4, 0.22, text, size=9.5, color=INK)
    if note:
        txt(s, x + 0.32, y + 0.22, COLW - 0.4, 0.2, note, size=8, color=AMBER)

VY = HY + 0.38 + 3 * 0.54 + 0.12
hline(s, ML, VY, CW, RULE, 0.75)
for dx, w, lbl, fc, tc in [(0.00, 1.90, "Comparable", GREEN_LT, GREEN),
                           (2.05, 2.55, "Conditionally comparable", AMBER_LT, AMBER),
                           (4.75, 1.90, "Not comparable", RED_LT, RED)]:
    rect(s, ML + dx, VY + 0.16, w, 0.32, fill=fc)
    txt(s, ML + dx, VY + 0.235, w, 0.2, lbl, size=9.5, font=F_SB, color=tc,
        align=PP_ALIGN.CENTER)
txt(s, ML + 7.05, VY + 0.14, CW - 7.05, 0.42,
    [[("Where the methodological gap could plausibly exceed the measured difference, ", {"color": INK2}),
      ("report no difference.", {"font": F_SB, "color": INK})]], size=9, spacing=1.24)
txt(s, ML, VY + 0.60, CW, 0.2,
    "The prototype does not yet demonstrate compliance with these conditions — page 12 names the fields that would close it.",
    size=8.5, color=INK3)

# ══════════════════════════════════════════════════════════ PAGE 7
# Mode: evidence chain. One value, three states.
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=PAPER)
page_frame(
    s, 1, ACTS[1],
    [[("A number without its evidence is not a measurement — and a", {})],
     [("missing number is never a zero", {})]],
    "Evidence is four separate things, and none of them may be folded into the performance number. "
    "If absence defaults to zero, the least transparent supplier becomes the best performer.",
    "P3 PACT v3 reliability KPIs, incl. share of primary data  ·  A2 ISO 14044 temporal, geographical and technological representativeness  ·  "
    "P1 GHG Protocol Product Standard §A.1, conditions 4 and 6  ·  P2 JRC132067, mixed primary-foreground / secondary-background inventory structure  ·  E2.  "
    "The four-way split of evidence is a source-supported interpretation.  The value shown is illustrative and is not taken from the dataset.", 7)

PW7, PG7, PH7 = 3.72, 0.3065, 3.10
VAL = "2.4 kg CO₂e / kg"
for i, (state, fc, tc) in enumerate([("BARE", AMBER, INK), ("QUALIFIED", GREEN, PAPER),
                                     ("ABSENT", RULE, INK)]):
    x = ML + i * (PW7 + PG7)
    rect(s, x, Y_EX, PW7, PH7, fill=OFFWHITE)
    rect(s, x, Y_EX, PW7, 0.32, fill=fc)
    txt(s, x + 0.26, Y_EX + 0.065, PW7 - 0.5, 0.2, state, size=9.5, font=F_SB, color=tc)

    if i < 2:
        txt(s, x + 0.26, Y_EX + 0.52, PW7 - 0.5, 0.4, VAL, size=19, font=F_L, color=INK)
    else:
        rect(s, x + 0.26, Y_EX + 0.52, 1.70, 0.42, fill=PAPER, line=RULE, lw=0.75)
        txt(s, x + 0.26, Y_EX + 0.615, 1.70, 0.22, "—", size=13, font=F_L,
            color=INK3, align=PP_ALIGN.CENTER)
        txt(s, x + 2.06, Y_EX + 0.60, 1.5, 0.24, "no figure reported", size=9, color=INK3)
    hline(s, x + 0.26, Y_EX + 1.08, PW7 - 0.52, RULE, 0.75)

    if i == 0:
        for j, a in enumerate(["Provenance", "Data quality", "Verification", "Uncertainty"]):
            txt(s, x + 0.26, Y_EX + 1.24 + j * 0.30, PW7 - 0.52, 0.22,
                [[(a, {"color": INK2}), ("      ?", {"font": F_SB, "color": AMBER})]], size=9.5)
        txt(s, x + 0.26, Y_EX + 2.62, PW7 - 0.52, 0.24,
            "A number. Not a measurement.", size=10, font=F_SB, color=INK)
    elif i == 1:
        for j, (a, v) in enumerate([
            ("Provenance", "supplier-specific; primary-data share"),
            ("Data quality", "temporal, geographical, technological"),
            ("Verification", "third-party assurance performed"),
            ("Uncertainty", "assessed and reported with the value"),
        ]):
            y = Y_EX + 1.24 + j * 0.30
            txt(s, x + 0.26, y, PW7 - 0.52, 0.22,
                [[(a, {"font": F_SB, "color": INK}), ("   " + v, {"color": INK2})]],
                size=8.5)
        txt(s, x + 0.26, Y_EX + 2.62, PW7 - 0.52, 0.24,
            "The same number, now usable.", size=10, font=F_SB, color=INK)
    else:
        txt(s, x + 0.26, Y_EX + 1.24, 0.5, 0.34, "0", size=19, font=F_L, color=RED)
        cz = s.shapes.add_connector(1, In(x + 0.22), In(Y_EX + 1.36),
                                    In(x + 0.60), In(Y_EX + 1.36))
        cz.line.color.rgb = RED; cz.line.width = Pt(2.0)
        txt(s, x + 0.68, Y_EX + 1.30, PW7 - 0.94, 0.22, "never", size=10, font=F_SB, color=RED)
        txt(s, x + 0.26, Y_EX + 1.74, PW7 - 0.52, 0.6,
            "If absence defaults to zero, the least transparent supplier becomes the "
            "best performer.", size=9, color=INK2, spacing=1.22)
        txt(s, x + 0.26, Y_EX + 2.62, PW7 - 0.52, 0.24,
            "Missing stays missing.", size=10, font=F_SB, color=INK)

SY7 = Y_EX + 3.26
hline(s, ML, SY7, CW, RULE, 0.75)
txt(s, ML, SY7 + 0.14, 6.55, 0.5,
    [[("Supplier-specific is not the same as primary.  ", {"font": F_SB, "color": INK}),
      ("A supplier-specific footprint normally combines primary activity data with secondary "
       "background factors — which is why PACT reports a share of primary data, not a flag.",
       {"color": INK2})]], size=9, spacing=1.24)
txt(s, ML + 7.05, SY7 + 0.14, CW - 7.05, 0.5,
    [[("Carried into the prototype.  ", {"font": F_SB, "color": INK}),
      ("It implements provenance and temporal quality, and enforces the never-zero rule at "
       "four independent layers — page 8.", {"color": INK2})]], size=9, spacing=1.24)

HEATMAP = os.path.join(HERE, "assets",
                       "page09_portfolio_heatmap.png")


def synth_chip(sl, x, y, w=2.62):
    """Mandatory from page 8: the synthetic-data boundary, visible on the page."""
    rect(sl, x, y, w, 0.22, fill=None, line=RULE, lw=0.75)
    txt(sl, x + 0.12, y + 0.045, w - 0.24, 0.18, "Synthetic demonstration data",
        size=7.5, font=F_SB, color=INK3, caps=True)


# ══════════════════════════════════════════════════════════ PAGE 8
# Mode: prototype scope. What the demonstrator is, and what the totals mean.
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=PAPER)
page_frame(
    s, 2, ACTS[2],
    [[("A working demonstrator on synthetic data — where the portfolio", {})],
     [("total is five times smaller than it first appears", {})]],
    "Five benchmark cases, twenty-five material alternatives, one incumbent in each. The demonstrator "
    "carries the decision structure end to end — and its first lesson is an arithmetic one.",
    "SYNTHETIC DEMONSTRATION DATA — no supplier, company or market data is used or implied.  ·  Figures computed from "
    "data/processed/consolidated_material_benchmark.csv through the 21 audited semantic measures; percentage deltas are "
    "recomputed from absolute values per audit item A3.", 8)
synth_chip(s, ML + 3.95, Y_EYE - 0.015)

# ---- left: the portfolio-total correction
LW8 = 6.90
txt(s, ML, Y_EX, LW8, 0.2, "Portfolio position", size=8.5, font=F_SB, color=INK3, caps=True)
hline(s, ML, Y_EX + 0.26, LW8, INK, 1.0)

txt(s, ML, Y_EX + 0.36, LW8, 0.2, "All 25 options, summed", size=9, color=INK3)
txt(s, ML, Y_EX + 0.54, 4.2, 0.5, "4,073,400 €", size=26, font=F_L, color=INK3)
c = s.shapes.add_connector(1, In(ML - 0.04), In(Y_EX + 0.79), In(ML + 2.70), In(Y_EX + 0.79))
c.line.color.rgb = RED; c.line.width = Pt(2.25)
txt(s, ML, Y_EX + 1.12, LW8 - 0.3, 0.5,
    [[("Every alternative is a scenario for the same demand. ", {"color": INK2}),
      ("Summing them counts that demand five times over.", {"font": F_SB, "color": INK})]],
    size=9.5, spacing=1.24)
hline(s, ML, Y_EX + 1.64, LW8, RULE, 0.75)

txt(s, ML, Y_EX + 1.76, LW8, 0.2, "The incumbent portfolio — the position actually held",
    size=9, color=INK3)
txt(s, ML, Y_EX + 1.94, 4.6, 0.58, "820,000 €", size=30, font=F_L, color=INK)
txt(s, ML, Y_EX + 2.58, LW8, 0.2,
    [[("480,000 kg", {"font": F_SB, "color": INK}), ("  annual volume     ", {"color": INK3}),
      ("605,000 kg", {"font": F_SB, "color": INK}), ("  annual CO₂e", {"color": INK3})]], size=9)

# ---- right: what the demonstrator contains
RX8, RW8 = ML + 7.25, 4.52
txt(s, RX8, Y_EX, RW8, 0.2, "The demonstrator", size=8.5, font=F_SB, color=INK3, caps=True)
hline(s, RX8, Y_EX + 0.26, RW8, INK, 1.0)
for j, (cid, grp) in enumerate([("CASE-A", "Base Oils"), ("CASE-B", "Solvents"),
                                ("CASE-C", "Surfactants"), ("CASE-D", "Polymer Additives"),
                                ("CASE-E", "Resins / Binders")]):
    y = Y_EX + 0.36 + j * 0.28
    txt(s, RX8, y, 1.0, 0.2, cid, size=9, font=F_SB, color=INK)
    txt(s, RX8 + 1.02, y, 2.0, 0.2, grp, size=9, color=INK2)
    txt(s, RX8 + 2.95, y, RW8 - 2.95, 0.2, "1 incumbent · 4 alternatives", size=8.5, color=INK3)

txt(s, RX8, Y_EX + 1.92, RW8, 0.2, "Coverage across all 25", size=8.5, font=F_SB,
    color=INK3, caps=True)
hline(s, RX8, Y_EX + 2.18, RW8, INK, 1.0)
for j, (v, lbl, sub) in enumerate([
        ("23 / 25", "technically eligible", "92.0%"),
        ("24 / 25", "PCF available", "96.0%"),
        ("96.8%", "of spend covered by a PCF", "")]):
    y = Y_EX + 2.28 + j * 0.26
    txt(s, RX8, y, 1.0, 0.2, v, size=9.5, font=F_SB, color=INK)
    txt(s, RX8 + 1.02, y, 2.6, 0.2, lbl, size=9, color=INK2)
    if sub:
        txt(s, RX8 + 3.62, y, RW8 - 3.62, 0.2, sub, size=9, color=INK3,
            align=PP_ALIGN.RIGHT)

# ---- bottom: the boundary of the claim
hline(s, ML, Y_EX + 3.24, CW, RULE, 0.75)
for dx, head, body in [
    (0.0, "What the demonstrator shows",
     "The decision structure end to end: gates that screen without deleting, evidence carried "
     "beside performance, and every missing value preserved as missing."),
    (6.05, "What it does not show",
     "A validated supplier comparison. No conformance to the six §A.1 conditions is claimed, and no "
     "real supplier, market or company data is used — page 12 states exactly what is absent."),
]:
    txt(s, ML + dx, Y_EX + 3.36, 5.6, 0.2, head, size=9.5, font=F_SB, color=INK)
    txt(s, ML + dx, Y_EX + 3.58, 5.6, 0.44, body, size=8.5, color=INK2, spacing=1.24)

# ══════════════════════════════════════════════════════════ PAGE 9
# Mode: flagship. Compressed header so the exhibit takes the page.
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=PAPER)
txt(s, ML, 0.30, 7.4, 0.22, ACTS[2], size=8.5, font=F_SB, color=INK3, caps=True)
sx = W - MR - (5 * 0.30 + 4 * 0.07)
for i in range(5):
    rect(s, sx + i * 0.37, 0.345, 0.30, 0.075, fill=BLUE if i == 2 else RULE_LT)
synth_chip(s, ML + 3.95, 0.285)
txt(s, ML, 0.56, CW - 0.4, 0.95,
    [[("All 25 alternatives on one surface: cost and carbon as the", {})],
     [("analysis, technical and evidence status as exceptions", {})]],
    size=22, font=F_SB, color=INK, spacing=1.06)
hline(s, ML, 1.34, CW, RULE, 0.75)

for dx, fc, lbl in [(0.00, C(0x3E, 0x8F, 0x63), "improvement against the incumbent of the same case"),
                    (3.55, C(0xA3, 0x2E, 0x22), "deterioration, or a technical blocker"),
                    (6.45, AMBER, "weak, outdated or missing PCF evidence"),
                    (9.35, PAPER, "incumbent, or normal state")]:
    rect(s, ML + dx, 1.475, 0.17, 0.15, fill=fc, line=RULE if fc == PAPER else None, lw=0.75)
    txt(s, ML + dx + 0.23, 1.45, 2.90, 0.2, lbl, size=8.5, color=INK2)

s.shapes.add_picture(HEATMAP, In(ML), In(1.68), In(CW), In(5.00))

txt(s, ML, Y_FOOT, CW - 1.0, 0.42,
    "SYNTHETIC DEMONSTRATION DATA.  Report-native re-render of the Power BI Python visual "
    "(powerbi/…/visuals/99aabbccddeeff001122) — plotting logic reused verbatim; only the figure size and the "
    "chrome now supplied by this page differ.  ·  Deltas are computed against the incumbent of each case from "
    "absolute spend and CO₂e.  ·  n / a is unknown and is never read as zero  ·  ▸ marks the incumbent  ·  "
    "each case is an independent decision  ·  no score, no ranking, no weighting.",
    size=7.5, color=INK3, spacing=1.22)
txt(s, W - MR - 0.6, Y_FOOT, 0.6, 0.22, "9", size=9, font=F_SB, color=INK3, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════ PAGE 10
# Mode: worked case — eligibility.
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=PAPER)
page_frame(
    s, 2, ACTS[2] + "   ·   CASE-D  POLYMER ADDITIVES",
    [[("The two strongest alternatives in this case", {})],
     [("are the two you cannot currently buy", {})]],
    "On cost and carbon the answer looks obvious. On eligibility, two of the five are not purchasable — "
    "and no performance number changes that.",
    "SYNTHETIC DEMONSTRATION DATA.  CASE-D, all five materials, from the governed dataset via the audited measures.  ·  "
    "Eligibility is the prototype gate (technically_eligible, technical_approval_status); it is shown, never used to filter "
    "rows away — BP-1.  ·  Screening removes an option from the comparison, not from the record.", 10)
synth_chip(s, ML + 5.55, Y_EYE - 0.015)

# ---- scatter, drawn natively so it stays editable
PX, PY, PWD, PHT = ML + 0.62, Y_EX + 0.34, 5.55, 2.52
XMIN, XMAX, YMIN, YMAX = -18.0, 7.0, -46.0, 14.0
fx = lambda v: PX + (v - XMIN) / (XMAX - XMIN) * PWD
fy = lambda v: PY + (YMAX - v) / (YMAX - YMIN) * PHT
rect(s, PX, PY, PWD, PHT, fill=OFFWHITE)
hline(s, PX, fy(0), PWD, INK3, 0.75)
c = s.shapes.add_connector(1, In(fx(0)), In(PY), In(fx(0)), In(PY + PHT))
c.line.color.rgb = INK3; c.line.width = Pt(0.75)
txt(s, PX + 0.10, fy(0) + 0.10, 2.0, 0.2, "better on both", size=8.5, font=F_SB, color=INK3)
txt(s, PX, PY + PHT + 0.10, PWD, 0.2, "Cost delta % vs incumbent      cheaper ←      → dearer",
    size=8.5, color=INK3, align=PP_ALIGN.CENTER)
txt(s, ML, PY - 0.24, 3.4, 0.2, "CO₂e delta % vs incumbent      lower carbon ↓",
    size=8.5, color=INK3)

for mid, cd, gd, fc, lab_dx, lab_dy in [
        ("PA-4001", 0.00, 0.00, INK, 0.14, -0.09),
        ("PA-4002", -14.29, -38.89, RED, 0.14, -0.09),
        ("PA-4003", -3.57, -5.56, BLUE, 0.14, -0.09),
        ("PA-4004", -7.14, -16.67, RED, 0.14, -0.09),
        ("PA-4005", 3.57, 8.33, BLUE, -1.14, -0.09)]:
    px, py = fx(cd), fy(gd)
    rect(s, px - 0.075, py - 0.075, 0.15, 0.15, fill=fc, shape=MSO_SHAPE.OVAL)
    txt(s, px + lab_dx, py + lab_dy, 0.95, 0.2, mid, size=8.5, font=F_SB,
        color=INK, align=PP_ALIGN.RIGHT if lab_dx < 0 else PP_ALIGN.LEFT)
txt(s, fx(-14.29) + 0.14, fy(-38.89) + 0.09, 1.9, 0.2, "blocked", size=8, font=F_SB, color=RED)
txt(s, fx(-7.14) + 0.14, fy(-16.67) + 0.09, 1.9, 0.2, "blocked", size=8, font=F_SB, color=RED)

# ---- right: the eligibility record
RX10, RW10 = ML + 6.95, 4.82
txt(s, RX10, Y_EX, RW10, 0.2, "Every option stays on the record", size=8.5, font=F_SB,
    color=INK3, caps=True)
hline(s, RX10, Y_EX + 0.26, RW10, INK, 1.0)
for j, h, w_, al in [(0, "", 1.30, PP_ALIGN.LEFT), (1, "Cost Δ", 0.92, PP_ALIGN.RIGHT),
                     (2, "CO₂e Δ", 0.92, PP_ALIGN.RIGHT), (3, "Technical status", 1.60, PP_ALIGN.LEFT)]:
    xo = [0.0, 1.30, 2.30, 3.35][j]
    txt(s, RX10 + xo, Y_EX + 0.32, w_, 0.18, h, size=7.5, font=F_SB, color=INK3,
        caps=True, align=al)
for j, (mid, cd, gd, st, blocked) in enumerate([
        ("PA-4001", "0.00%", "0.00%", "Incumbent", None),
        ("PA-4002", "−14.29%", "−38.89%", "Not Approved", True),
        ("PA-4003", "−3.57%", "−5.56%", "Approved", False),
        ("PA-4004", "−7.14%", "−16.67%", "Under Qualification", True),
        ("PA-4005", "+3.57%", "+8.33%", "Approved", False)]):
    y = Y_EX + 0.56 + j * 0.30
    if blocked:
        rect(s, RX10, y - 0.035, RW10, 0.29, fill=RED_LT)
    txt(s, RX10 + 0.06, y, 1.24, 0.2, mid, size=9, font=F_SB, color=INK)
    txt(s, RX10 + 1.30, y, 0.92, 0.2, cd, size=9, color=INK2, align=PP_ALIGN.RIGHT)
    txt(s, RX10 + 2.30, y, 0.92, 0.2, gd, size=9, color=INK2, align=PP_ALIGN.RIGHT)
    txt(s, RX10 + 3.35, y, 1.45, 0.2, st, size=9,
        font=F_SB if blocked else F_R, color=RED if blocked else INK2)

rect(s, RX10, Y_EX + 2.20, RW10, 1.16, fill=OFFWHITE)
rect(s, RX10, Y_EX + 2.20, 0.035, 1.16, fill=BLUE)
txt(s, RX10 + 0.26, Y_EX + 2.34, RW10 - 0.5, 0.84,
    [[("The only eligible improver is PA-4003 — ", {"font": F_SB, "color": INK}),
      ("−3.57% cost and −5.56% CO₂e. That, not −14.29%, is the decision actually on the table.",
       {"color": INK2})],
     [("", {})],
     [("PA-4004 is under qualification and may return. Both blocked options stay visible, with the "
       "reason attached.", {"color": INK3})]], size=9, spacing=1.22, space_after=2)

txt(s, ML, Y_EX + 3.52, CW, 0.22,
    [[("Performance cannot compensate for failed eligibility.", {"font": F_SB, "color": INK}),
      ("   A gate is passed or it is not; it is never traded off against a better number — page 4.",
       {"color": INK3})]], size=9)

# ══════════════════════════════════════════════════════════ PAGE 11
# Mode: worked case — evidence.
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=PAPER)
page_frame(
    s, 2, ACTS[2] + "   ·   CASE-C  SURFACTANTS",
    [[("The best carbon result rests on the weakest evidence — and", {})],
     [("one material has no carbon figure at all", {})]],
    "All five are technically approved, so eligibility decides nothing here. Evidence decides everything — "
    "and the two are read in different columns.",
    "SYNTHETIC DEMONSTRATION DATA.  CASE-C, all five materials, from the governed dataset via the audited measures.  ·  "
    "BP-2 an industry-average factor cannot distinguish two suppliers within the same industry  ·  BP-3 evidence quality is "
    "carried beside the number  ·  BP-4 missing is never zero  ·  P3 PACT v3  ·  A2 ISO 14044.", 11)
synth_chip(s, ML + 5.05, Y_EYE - 0.015)

COLS11 = [("Material", 2.30, PP_ALIGN.LEFT), ("Cost Δ %", 1.15, PP_ALIGN.RIGHT),
          ("CO₂e Δ %", 1.20, PP_ALIGN.RIGHT), ("PCF kg CO₂e / kg", 1.55, PP_ALIGN.RIGHT),
          ("PCF data type", 2.85, PP_ALIGN.LEFT), ("Year", 0.75, PP_ALIGN.RIGHT),
          ("Confidence", 1.40, PP_ALIGN.LEFT)]
xo, XO11 = 0.0, []
for _, w_, _a in COLS11:
    XO11.append(xo); xo += w_ + 0.06
XO11[6] += 0.20                      # breathing room before Confidence
for j, (h, w_, al) in enumerate(COLS11):
    txt(s, ML + XO11[j], Y_EX, w_, 0.18, h, size=7.5, font=F_SB, color=INK3, caps=True, align=al)
hline(s, ML, Y_EX + 0.24, CW, INK, 1.0)

ROWS11 = [
    ("SF-3001", "0.00%", "0.00%", "2.10", "Supplier-specific (primary)", "2023", "High confidence", None),
    ("SF-3002", "+6.25%", "−54.76%", "0.95", "Industry-average (secondary)", "2019", "Low confidence", "amber"),
    ("SF-3003", "−8.33%", "", "", "no PCF reported", "", "No data", "blank"),
    ("SF-3004", "+10.42%", "−23.81%", "1.60", "Supplier-specific (primary)", "2023", "High confidence", None),
    ("SF-3005", "−2.08%", "+7.14%", "2.25", "Supplier-specific (primary)", "2023", "High confidence", None),
]
for i, (mid, cd, gd, pv, dt, yr, cf, mark) in enumerate(ROWS11):
    y = Y_EX + 0.34 + i * 0.34
    if mark == "amber":
        rect(s, ML, y - 0.045, CW, 0.33, fill=AMBER_LT)
    elif i % 2 == 1:
        rect(s, ML, y - 0.045, CW, 0.33, fill=OFFWHITE)
    vals = [mid, cd, gd, pv, dt, yr, cf]
    for j, (v, (_h, w_, al)) in enumerate(zip(vals, COLS11)):
        if v == "" and j in (2, 3, 5):
            rect(s, ML + XO11[j] + w_ - 0.42, y + 0.10, 0.34, 0.02, fill=RULE)
            continue
        bold = (j == 0) or (mark == "amber" and j in (2, 4, 5, 6)) or (mark == "blank" and j in (4, 6))
        col = INK if j == 0 else (INK2 if mark is None else INK)
        txt(s, ML + XO11[j], y, w_, 0.22, v, size=9,
            font=F_SB if bold else F_R, color=col, align=al)

hline(s, ML, Y_EX + 2.10, CW, RULE, 0.75)
for dx, bar, head, body in [
    (0.0, AMBER, "SF-3002 — the strongest carbon result in the case",
     "−54.76% CO₂e, carried by the weakest evidence on the page: an industry-average factor from 2019, "
     "low confidence. An industry-average factor cannot distinguish two suppliers within the same industry. "
     "The number is not wrong — it is not yet a supplier comparison."),
    (6.05, INK3, "SF-3003 — no carbon figure at all",
     "−8.33% on cost, and blank on carbon everywhere: in the measure, in the heatmap, in the case total. "
     "It is not zero, it is not a worse performer, and it is not dropped from the cost comparison. "
     "It is unknown, and it stays unknown."),
]:
    rect(s, ML + dx, Y_EX + 2.26, 5.72, 1.16, fill=OFFWHITE)
    rect(s, ML + dx, Y_EX + 2.26, 0.035, 1.16, fill=bar)
    txt(s, ML + dx + 0.26, Y_EX + 2.40, 5.32, 0.2, head, size=9.5, font=F_SB, color=INK)
    txt(s, ML + dx + 0.26, Y_EX + 2.64, 5.32, 0.7, body, size=8.5, color=INK2, spacing=1.24)

txt(s, ML, Y_EX + 3.56, CW, 0.22,
    [[("The result and the confidence in the result are two columns.", {"font": F_SB, "color": INK}),
      ("   Blending them into one figure destroys both — page 7.", {"color": INK3})]], size=9)

# ══════════════════════════════════════════════════════════ PAGE 12
# Mode: honest scope boundary, measured against page 6.
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=PAPER)
page_frame(
    s, 2, ACTS[2],
    [[("The architecture is right; the comparison is not yet demonstrably", {})],
     [("valid — and we know which fields would close it", {})]],
    "Page 6 set six conditions. Measured against them the demonstrator is honest about what it cannot yet "
    "support — which is exactly what makes it usable as a specification.",
    "SYNTHETIC DEMONSTRATION DATA.  Conditions verbatim from the GHG Protocol Product Standard, Appendix A §A.1.  ·  "
    "Prototype status assessed against the fields present in the semantic model.  ·  "
    "No condition is fully satisfied today, so no status on this page is green.  ·  The closing fields are the "
    "prototype v0.2 candidate list.", 12)
synth_chip(s, ML + 3.95, Y_EYE - 0.015)

C12 = [("", 0.30), ("Condition — GHG Protocol §A.1", 3.30), ("In the prototype today", 3.05),
       ("Status", 1.42), ("The field that would close it", 3.30)]
xo, XO12 = 0.0, []
for _, w_ in C12:
    XO12.append(xo); xo += w_ + 0.09
for j, (h, w_) in enumerate(C12):
    txt(s, ML + XO12[j], Y_EX, w_, 0.18, h, size=7.5, font=F_SB, color=INK3, caps=True)
hline(s, ML, Y_EX + 0.24, CW, INK, 1.0)

ROWS12 = [
    ("1", "The unit of analysis is identical", "A per-kg declared unit only",
     "partial", "Functional or declared unit, with the reference flow"),
    ("2", "System and temporal boundaries are equivalent", "No boundary or data-period field",
     "not shown", "Cradle-to-gate boundary flag; data period"),
    ("3", "The same allocation methods are used", "No rulebook or version recorded",
     "not shown", "Rulebook and version — TfS, PACT, Catena-X"),
    ("4", "Data types, quality and uncertainty are reported", "pcf_data_type, pcf_data_quality_tier",
     "partial", "Primary-data share; uncertainty range"),
    ("5", "Temporal and geographical representativeness", "pcf_reference_year — temporal only",
     "partial", "Geographical representativeness"),
    ("6", "Third-party assurance has been performed", "No verification field at all",
     "not shown", "Verification status and assurance level"),
]
for i, (num, cond, today, st, fld) in enumerate(ROWS12):
    y = Y_EX + 0.34 + i * 0.50
    if i % 2 == 1:
        rect(s, ML, y - 0.06, CW, 0.48, fill=OFFWHITE)
    rect(s, ML + XO12[0], y + 0.015, 0.22, 0.20, fill=BLUE)
    txt(s, ML + XO12[0], y + 0.038, 0.22, 0.18, num, size=8, font=F_SB, color=PAPER,
        align=PP_ALIGN.CENTER)
    txt(s, ML + XO12[1], y, C12[1][1], 0.4, cond, size=9.5, color=INK, spacing=1.16)
    txt(s, ML + XO12[2], y + 0.015, C12[2][1], 0.4, today, size=8.5, color=INK2, spacing=1.16)
    fc, tc = (AMBER_LT, AMBER) if st == "partial" else (RULE_LT, INK3)
    rect(s, ML + XO12[3], y + 0.005, 1.18, 0.24, fill=fc)
    txt(s, ML + XO12[3], y + 0.045, 1.18, 0.2, st, size=8.5, font=F_SB, color=tc,
        align=PP_ALIGN.CENTER)
    txt(s, ML + XO12[4], y + 0.015, C12[4][1], 0.4, fld, size=8.5, color=INK2, spacing=1.16)

hline(s, ML, Y_EX + 3.36, CW, RULE, 0.75)
txt(s, ML, Y_EX + 3.48, 6.6, 0.6,
    [[("What the demonstrator does establish", {"font": F_SB, "color": INK})],
     [("Gates that screen without deleting, evidence carried beside performance, missing values "
       "preserved, and each case kept as a separate decision.", {"color": INK2})]],
    size=8.5, spacing=1.24)
txt(s, ML + 6.85, Y_EX + 3.48, CW - 6.85, 0.6,
    [[("Closing the six conditions is a data task, not a redesign", {"font": F_SB, "color": INK})],
     [("Every field in the right-hand column is collected upstream, in the specification and the RFQ. "
       "That scope is page 16.", {"color": INK2})]],
    size=8.5, spacing=1.24)

def ev_marker(sl, x, y, kind):
    """Page 5's three-state evidence marker, reused unchanged."""
    if kind == 2:
        rect(sl, x, y, 0.10, 0.10, fill=INK)
    elif kind == 1:
        rect(sl, x, y, 0.10, 0.10, fill=INK3)
    else:
        rect(sl, x, y, 0.10, 0.10, fill=PAPER, line=INK3, lw=0.75)


def ev_legend(sl, y):
    for dx, lbl, kind in [(0.00, "strongly evidenced — primary standard or EU authority", 2),
                          (3.60, "supported — cited, but narrower", 1),
                          (6.10, "project recommendation — plausible, not externally evidenced", 0)]:
        ev_marker(sl, ML + dx, y + 0.035, kind)
        txt(sl, ML + dx + 0.20, y, 3.5, 0.2, lbl, size=8.5, color=INK3)


# ══════════════════════════════════════════════════════════ PAGE 13
# Mode: bounded case evidence. Where circularity loses, and to what.
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=PAPER)
page_frame(
    s, 3, ACTS[3],
    [[("More circular is not automatically more sustainable — an EU", {})],
     [("assessment shows exactly where recovery loses", {})]],
    "Across 27 recycling pathways and 14 impact categories, recycling beat energy recovery on climate "
    "in every pathway analysed. On five other categories, energy-intensive routes lose.",
    "European Commission JRC, Environmental and economic assessment of plastic waste recycling, JRC132067 (2023) — executive summary read in "
    "the original; quoted wording retained.  ·  Supporting: C1 on recycling energy intensity.  ·  FACT, within the stated scope.  ·  "
    "The prototype does not measure circularity; nothing on this page comes from it.", 13)

# mandatory scope label — on the page, never in a footnote
rect(s, ML, Y_EX, CW, 0.46, fill=AMBER_LT)
rect(s, ML, Y_EX, 0.035, 0.46, fill=AMBER)
txt(s, ML + 0.26, Y_EX + 0.095, CW - 0.5, 0.2,
    [[("Scope   ", {"font": F_SB, "color": INK, "size": 8}),
      ("This is a plastics waste-treatment study. It demonstrates the mechanism — burden shifting between "
       "impact categories — and is not evidence about base oils, solvents, surfactants, polymer additives or resins.",
       {"color": INK})]], size=9)

# upper register — climate
rect(s, ML, Y_EX + 0.58, CW, 0.56, fill=GREEN_LT)
rect(s, ML, Y_EX + 0.58, 0.035, 0.56, fill=GREEN)
txt(s, ML + 0.26, Y_EX + 0.685, 2.0, 0.2, "Climate", size=10, font=F_SB, color=GREEN, caps=True)
txt(s, ML + 2.30, Y_EX + 0.68, CW - 2.6, 0.4,
    [[("Recycling beat energy recovery in every one of the 27 pathways analysed.", {"font": F_SB, "color": INK}),
      ("   This is the finding a circularity target is usually built on.", {"color": INK2})]],
    size=9.5, spacing=1.22)

# lower register — the five categories where the direction reverses
txt(s, ML, Y_EX + 1.32, 8.5, 0.2,
    "Where energy recovery can perform better — five of the fourteen categories",
    size=8.5, font=F_SB, color=INK3, caps=True)
hline(s, ML, Y_EX + 1.56, CW, INK, 1.0)
CHW, CHG = (CW - 4 * 0.13) / 5, 0.13
for j, cat in enumerate(["Acidification", "Particulate matter", "Ionising radiation",
                         "Human toxicity, non-cancer", "Eutrophication"]):
    x = ML + j * (CHW + CHG)
    rect(s, x, Y_EX + 1.66, CHW, 0.34, fill=RED_LT)
    txt(s, x, Y_EX + 1.735, CHW, 0.2, cat, size=9, font=F_SB, color=RED, align=PP_ALIGN.CENTER)

txt(s, ML, Y_EX + 2.14, CW, 0.2,
    [[("…for these energy-intensive recycling routes, named in the assessment", {"color": INK2})]],
    size=9)
for j, rt in enumerate(["PET alkaline hydrolysis", "PS mechanical", "MPO mechanical & pyrolysis",
                        "PE film mechanical & physical", "EPS physical"]):
    x = ML + j * (CHW + CHG)
    rect(s, x, Y_EX + 2.38, CHW, 0.30, fill=OFFWHITE)
    txt(s, x + 0.08, Y_EX + 2.445, CHW - 0.16, 0.2, rt, size=8.5, color=INK2,
        align=PP_ALIGN.CENTER)

hline(s, ML, Y_EX + 2.88, CW, RULE, 0.75)
txt(s, ML, Y_EX + 3.00, 5.6, 0.66,
    [[("Three further findings from the same assessment", {"font": F_SB, "color": INK})],
     [("Technical quality issues in the recovered material  ·  negative net incomes for methanolysis, "
       "pyrolysis and gasification  ·  and, across the whole study, “a clear ranking could not be established”.",
       {"color": INK2})]], size=8.5, spacing=1.26)
txt(s, ML + 6.05, Y_EX + 3.00, 5.7, 0.66,
    [[("And the answer has a date on it", {"font": F_SB, "color": INK})],
     [("The assessment attributes the result to the European energy mix, and states that as the mix gets "
       "cleaner the gap will widen further in favour of recycling.", {"color": INK2})]],
    size=8.5, spacing=1.26)

# ══════════════════════════════════════════════════════════ PAGE 14
# Mode: the trade-off map. Different tensions, different rules, no weights.
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=PAPER)
page_frame(
    s, 3, ACTS[3],
    [[("Every trade-off has its own handling rule — and none of", {})],
     [("them is a weight", {})]],
    "The tensions are structurally different. Some resolve at a gate, some are genuine two-axis trade-offs, "
    "one is about evidence rather than performance. Not one of them needs a weighting.",
    "T-1 to T-9: author synthesis, with the confidence recorded for each tension  ·  P2 JRC132067 “a clear ranking could not be established”  ·  "
    "P1 GHG Protocol Product Standard §1.7 on burden shifting  ·  V3 ISO 14046, V4 EU PEF  ·  BP-7.  The map is a "
    "source-supported interpretation; the individual evidence items are fact.  ·  Cost ↔ carbon quantitative figures in the register are "
    "steel-sector illustration only; no chemical-sector green-premium dataset was found in either research pass.", 14)

C14 = [("Tension", 2.85), ("How it shows up", 3.30), ("Evidence", 0.90), ("Handling rule", 4.35)]
xo, XO14 = 0.0, []
for _, w_ in C14:
    XO14.append(xo); xo += w_ + 0.11
for j, (h, w_) in enumerate(C14):
    txt(s, ML + XO14[j], Y_EX, w_, 0.18, h, size=7.5, font=F_SB, color=INK3, caps=True,
        align=PP_ALIGN.CENTER if j == 2 else PP_ALIGN.LEFT)
hline(s, ML, Y_EX + 0.24, CW, INK, 1.0)

ROWS14 = [
    ("Cost ↔ carbon", "The cheaper option carries more carbon, or the reverse", 0,
     "Two axes, plus a displayed cost-per-tonne ratio — never a score"),
    ("Circularity ↔ non-climate impacts", "Recovery wins on climate and loses elsewhere", 2,
     "Report as distinct outputs; any priority rule is stated per case"),
    ("Circularity ↔ technical performance", "Recycled content shifts the specification", 1,
     "Gate — recycled content is a candidate only if it passes equivalence"),
    ("Circularity ↔ cost", "The break-even arrives later than the target date", 1,
     "Treat the horizon as an explicit scenario, not an assumption"),
    ("Sustainability ↔ technical performance", "A carbon gain sits behind a failed approval", 2,
     "Gate. No carbon result buys past a technical blocker"),
    ("Sustainability ↔ supply security", "A carbon gain raises supplier concentration", 1,
     "State it beside the result; concentration needs awarded volumes"),
    ("Performance ↔ evidence quality", "The best number has the weakest provenance", 2,
     "Two columns. Never a quality-adjusted performance figure"),
    ("Carbon ↔ other environmental burdens", "A GHG-only view misses the shift", 2,
     "Disclose scope. Methods exist — the constraint is supplier data"),
]
for i, (ten, shows, ev, rule) in enumerate(ROWS14):
    y = Y_EX + 0.34 + i * 0.32
    if i % 2 == 1:
        rect(s, ML, y - 0.045, CW, 0.31, fill=OFFWHITE)
    txt(s, ML + XO14[0], y, C14[0][1], 0.3, ten, size=9, font=F_SB, color=INK, spacing=1.14)
    txt(s, ML + XO14[1], y + 0.01, C14[1][1], 0.3, shows, size=8.5, color=INK2, spacing=1.14)
    ev_marker(s, ML + XO14[2] + C14[2][1] / 2 - 0.05, y + 0.055, ev)
    txt(s, ML + XO14[3], y + 0.01, C14[3][1], 0.3, rule, size=8.5, color=INK2, spacing=1.14)
hline(s, ML, Y_EX + 2.92, CW, RULE, 0.75)

ev_legend(s, Y_EX + 3.00)

BY14 = Y_EX + 3.22
rect(s, ML, BY14, CW, 0.72, fill=INK)
txt(s, ML + 0.28, BY14 + 0.10, CW - 0.6, 0.58,
    [[("Deliberately rejected", {"font": F_SB, "color": PAPER})],
     [("A composite sustainability score  ·  arbitrary weighting  ·  cross-case ranking of "
       "non-substitutable materials. An EU authority reached the same conclusion in its own domain: "
       "“a clear ranking could not be established”.", {"color": C(0xC6, 0xD4, 0xE0)})]],
    size=9.5, spacing=1.30)

# ══════════════════════════════════════════════════════════ PAGE 15
# Mode: governance. Who owns what, and the artefact that carries it.
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=PAPER)
page_frame(
    s, 4, ACTS[4],
    [[("Four functions, one decision record, and a named owner", {})],
     [("for every threshold", {})]],
    "Every practice in this report fails at the same seam: gates, commercial comparison and evidence "
    "rules sit with different functions, and each keeps its own list.",
    "A1 ISO 20400 places sustainability criteria alongside price and quality in a single evaluation rather than in a parallel process  ·  "
    "BP-11: one comparison case, one decision record.  The ownership split and the record's field list are a PROJECT SYNTHESIS "
    "and a design recommendation: ISO 20400 supports the principle, it does not validate this artefact.", 15)

rect(s, ML + 4.35, Y_EYE - 0.015, 3.30, 0.22, fill=None, line=RULE, lw=0.75)
txt(s, ML + 4.47, Y_EYE + 0.03, 3.06, 0.18,
    [[("Project synthesis", {"font": F_SB, "color": INK3}),
      ("  ·  design recommendation", {"color": INK3})]], size=7.5, caps=True)

LW15 = 5.45
txt(s, ML, Y_EX, LW15, 0.2, "Who owns which decision", size=8.5, font=F_SB, color=INK3, caps=True)
hline(s, ML, Y_EX + 0.26, LW15, INK, 1.0)
for j, (fn, duty, key) in enumerate([
    ("Engineering", "Defines the function and validates technical equivalence", BLUE),
    ("Procurement", "Commercial terms, capacity evidence, sourcing scenarios", BLUE),
    ("Sustainability", "Accepted methods and evidence requirements", BLUE),
    ("Compliance", "Specific legal and policy gates; the due-diligence risk process", BLUE),
    ("Suppliers", "Provide declarations and evidence — never self-approve gate status", RULE),
    ("Business decision owner", "Accepts the residual trade-off, on the record", INK),
]):
    y = Y_EX + 0.38 + j * 0.50
    rect(s, ML, y, LW15, 0.44, fill=OFFWHITE)
    rect(s, ML, y, 0.035, 0.44, fill=key)
    txt(s, ML + 0.24, y + 0.045, LW15 - 0.45, 0.2, fn, size=9.5, font=F_SB, color=INK)
    txt(s, ML + 0.24, y + 0.235, LW15 - 0.45, 0.2, duty, size=8.5, color=INK2)

RX15, RW15 = ML + 5.85, 5.93
txt(s, RX15, Y_EX, RW15, 0.2, "One decision record per benchmark case", size=8.5, font=F_SB,
    color=INK3, caps=True)
hline(s, RX15, Y_EX + 0.26, RW15, INK, 1.0)
rect(s, RX15, Y_EX + 0.38, RW15, 2.14, fill=OFFWHITE)
for j, (fld, note) in enumerate([
    ("Incumbent and candidates", "the case, not the catalogue"),
    ("Gate status, with the reason", "blocked options stay visible"),
    ("Cost position  ·  carbon position", "two axes, never merged"),
    ("Evidence status per number", "provenance, quality, verification, uncertainty"),
    ("Comparability verdict", "comparable / conditional / not comparable"),
    ("The decision", "and what was given up"),
    ("Named owner", "a threshold without one is an opinion"),
    ("Review date", "the record expires"),
]):
    y = Y_EX + 0.50 + j * 0.25
    rect(s, RX15 + 0.24, y + 0.055, 0.075, 0.075, fill=BLUE)
    txt(s, RX15 + 0.44, y, 2.85, 0.2, fld, size=9, font=F_SB, color=INK)
    txt(s, RX15 + 3.34, y, RW15 - 3.5, 0.2, note, size=8.5, color=INK3)

rect(s, RX15, Y_EX + 2.64, RW15, 0.62, fill=None, line=RULE, lw=0.75)
txt(s, RX15 + 0.24, Y_EX + 2.75, RW15 - 0.48, 0.44,
    [[("A second output, not a footnote: the qualification pipeline. ", {"font": F_SB, "color": INK}),
      ("Which blocked candidates are worth qualifying, by when, and at whose cost.", {"color": INK2})]],
    size=8.5, spacing=1.26)

hline(s, ML, Y_EX + 3.42, CW, RULE, 0.75)
txt(s, ML, Y_EX + 3.54, CW, 0.22,
    [[("The prototype is not yet this artefact.", {"font": F_SB, "color": INK}),
      ("   Power BI page 2 is the closest thing that exists: it displays the decision context, but it "
       "does not capture the decision, the owner or the review date.", {"color": INK3})]], size=9)

# ══════════════════════════════════════════════════════════ PAGE 16
# Mode: closing roadmap. Build, build-if-material, and refuse.
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=PAPER)
page_frame(
    s, 4, ACTS[4],
    [[("Two lanes to build, one list to refuse — and the gaps we", {})],
     [("are not pretending to have closed", {})]],
    "Lane A makes the comparison already being made defensible. Lane B adds category modules where "
    "materiality justifies them. The third list is what must never be synthesised.",
    "Categories 1–4 and the two lanes: project synthesis  ·  limitations disclosed and not closed  ·  "
    "Lane A derives entirely from GHG Protocol Product Standard §A.1.  PROJECT SYNTHESIS, anchored on the frozen classification.", 16)

CW16, CG16 = (CW - 2 * 0.31) / 3, 0.31
COLS16 = [
    ("Build now — Lane A", "makes the existing comparison defensible", BLUE, PAPER, [
        "Functional / declared unit and dosage normalisation",
        "Methodology rulebook and version",
        "Boundary description or comparability flag",
        "Primary-data share (%)",
        "Verification status and scope",
        "Data / reference period",
        "Uncertainty, or an uncertainty description",
        "Comparable / conditional / not comparable verdict"],
     "All eight derive from GHG Protocol §A.1 — the conditions set out on page 12."),
    ("Build where material — Lane B", "only where the decision can change", BLUE_LT, INK, [
        "Recycled and renewable content, circularity",
        "Due diligence and human-rights performance",
        "Hazardous-substance legal status",
        "Criticality and supply resilience",
        "Water, land use, biodiversity, pollution, toxicity",
        "CBAM exposure where the product code is out of scope"],
     "Each gated by the materiality screen on page 3, with its own data owner and grain."),
    ("Do not build", "not until real data exists", INK, PAPER, [
        "Qualification and switching cost",
        "True TCO effects on yield, energy, quality, downtime",
        "Supplier capacity, ramp-up, allocation, recovery time",
        "Country / supplier concentration on awarded shares",
        "Toxicity, exposure and legal use status for real substances",
        "Recycled / renewable chain of custody and mass balance",
        "Responsible-sourcing findings and traceability depth"],
     "Synthesising any of these would demonstrate columns, not capability."),
]
for i, (head, sub, hfill, htxt, items, note) in enumerate(COLS16):
    x = ML + i * (CW16 + CG16)
    rect(s, x, Y_EX, CW16, 0.44, fill=hfill)
    txt(s, x + 0.18, Y_EX + 0.055, CW16 - 0.34, 0.2, head, size=10, font=F_SB, color=htxt)
    txt(s, x + 0.18, Y_EX + 0.255, CW16 - 0.34, 0.18, sub, size=8,
        color=C(0xC6, 0xD4, 0xE0) if htxt == PAPER else INK2)
    for j, it in enumerate(items):
        y = Y_EX + 0.56 + j * 0.27
        rect(s, x + 0.02, y + 0.055, 0.075, 0.075,
             fill=RULE if i == 2 else (BLUE if i == 0 else BLUE_LT))
        txt(s, x + 0.20, y, CW16 - 0.26, 0.24, it, size=8.5, color=INK2, spacing=1.16)
    txt(s, x + 0.02, Y_EX + 2.80, CW16 - 0.1, 0.44, note, size=8.5, font=F_SB,
        color=INK3, spacing=1.24)

hline(s, ML, Y_EX + 3.28, CW, RULE, 0.75)
txt(s, ML, Y_EX + 3.38, 5.65, 0.72,
    [[("Research limitations — disclosed, not closed", {"font": F_SB, "color": INK})],
     [("No industrial chemical-procurement case  ·  no empirical study of how far real supplier PCFs "
       "diverge when the six conditions are violated  ·  TfS v3 and Catena-X v4 not read in full  ·  no "
       "company-level resilience metric set  ·  no chemical-sector green-premium dataset.", {"color": INK2})]],
    size=8.5, spacing=1.24)
txt(s, ML + 6.05, Y_EX + 3.38, 5.7, 0.72,
    [[("Implementation prerequisites — not defects in the method", {"font": F_SB, "color": INK})],
     [("Every Lane A field is collected upstream, in the specification and the RFQ. Comparability is "
       "engineered there — it cannot be repaired downstream in the analysis.", {"color": INK2})]],
    size=8.5, spacing=1.24)

# ── document properties ───────────────────────────────────────────
# python-pptx inherits its template's author; set them so the file is
# attributed to the project author rather than to the library's.
cp = prs.core_properties
cp.author = "Csaba Bakay"
cp.last_modified_by = "Csaba Bakay"
cp.title = "Sustainable Raw Material Benchmarking"
cp.comments = ""

# ── save ──────────────────────────────────────────────────────────────────
os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print("saved:", OUT)
print("slides:", len(prs.slides.__iter__.__self__._sldIdLst))
